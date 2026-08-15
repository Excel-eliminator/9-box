import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import io
import os
import textwrap
import hashlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

st.set_page_config(page_title="Unified 9-Box Margin Analysis", layout="wide")


# ==========================================
# DATA LOADING FUNCTIONS
# ==========================================
# @st.cache_data -> DIHAPUS agar Streamlit selalu membaca data dari awal (100% fresh, no stubborn cache)
def load_9box_data(file_path):
    try:
        xls = pd.ExcelFile(file_path)
        df_local = pd.read_excel(xls, sheet_name='LOCAL')
        df_export = pd.read_excel(xls, sheet_name='EXPORT')

        df_local['Source_Sheet'] = 'LOCAL'
        df_export['Source_Sheet'] = 'EXPORT'

        df_combined = pd.concat([df_local, df_export], ignore_index=True)

        # Standardize numeric columns - MARGIN DIAMBIL MURNI DARI EXCEL TANPA DIUBAH!
        numeric_cols = [
            'Gross Sales (Current)', 'Return (Current)', 'COGS_Regular (Current)', 'Royalty (Current)',
            'Gross Profit (Current)', 'Gross Margin (Current)', 'Contribution Margin (Current)',
            'Gross Profit (%)', 'Gross Margin (%)', 'Contribution Margin (%)',
            'Qty (Current)', 'Gross Sales (Previous)', 'Qty (Previous)',
            'Qty Growth (%)', 'Amount FG', 'Amount Material'
        ]
        for col in numeric_cols:
            if col in df_combined.columns:
                df_combined[col] = pd.to_numeric(df_combined[col], errors='coerce').fillna(0)

        # Explicitly cast categorical/text fields to STRING to prevent Excel number conversion
        text_cols = ['SKU', 'Product Name', 'New Code', 'New Product Name', 'Remark', 'Status', 'Country']
        for col in text_cols:
            if col in df_combined.columns:
                df_combined[col] = df_combined[col].fillna('').astype(str)

        # --- FALLBACK CALCULATION JIKA PERSENTASE TIDAK ADA DI EXCEL ---
        if 'Qty Growth (%)' not in df_combined.columns:
            df_combined['Qty Growth (%)'] = 0.0

        if 'Gross Profit (%)' not in df_combined.columns and 'Gross Profit (Current)' in df_combined.columns and 'Gross Sales (Current)' in df_combined.columns:
            df_combined['Gross Profit (%)'] = np.where(df_combined['Gross Sales (Current)'] > 0, (df_combined['Gross Profit (Current)'] / df_combined['Gross Sales (Current)']) * 100, 0.0)

        if 'Gross Margin (%)' not in df_combined.columns and 'Gross Margin (Current)' in df_combined.columns and 'Gross Sales (Current)' in df_combined.columns:
            df_combined['Gross Margin (%)'] = np.where(df_combined['Gross Sales (Current)'] > 0, (df_combined['Gross Margin (Current)'] / df_combined['Gross Sales (Current)']) * 100, 0.0)

        if 'Contribution Margin (%)' not in df_combined.columns and 'Contribution Margin (Current)' in df_combined.columns and 'Gross Sales (Current)' in df_combined.columns:
            df_combined['Contribution Margin (%)'] = np.where(df_combined['Gross Sales (Current)'] > 0, (df_combined['Contribution Margin (Current)'] / df_combined['Gross Sales (Current)']) * 100, 0.0)

        return df_combined
    except Exception as e:
        st.error(f"Failed to load data: {e}")
        return pd.DataFrame()


# ==========================================
# 9-BOX SUMMARY GRID RENDERER
# ==========================================
def render_9box_summary_grid(df, margin_type, margin_val_col, y_low_thresh, y_high_thresh, x_axis_metric, x_low_thresh,
                             x_high_thresh):
    abbr = "GP" if margin_type == "Gross Profit" else ("GM" if margin_type == "Gross Margin" else "CM")

    def get_box_html(b_name, bg_color):
        sub = df[df['Dynamic 9-Box Category'].str.startswith(b_name)]
        cnt = len(sub)
        sales = sub['Gross Sales (Current)'].sum()
        margin = sub[margin_val_col].sum()
        qty = sub['Qty (Current)'].sum() if 'Qty (Current)' in sub.columns else 0

        pct = (margin / sales * 100) if sales > 0 else 0
        margin_bn = margin / 1e9
        sales_bn = sales / 1e9

        return f'<div class="box {bg_color}"><div class="box-title">{b_name}</div><div class="box-main">{cnt:,} <span class="box-sub">SKUs</span></div><div class="box-sub-small">IDR {margin_bn:,.0f} Bn margin &middot; {pct:.0f}% {abbr}<br>Sales: IDR {sales_bn:,.0f} Bn &middot; Qty: {qty:,.0f}</div></div>'

    def format_x_val(val):
        if x_axis_metric == "Qty Growth (%)":
            return f"{val:.1f}%"
        else:
            if val >= 1e9 or val <= -1e9:
                return f"Rp{val / 1e9:,.1f}M"
            else:
                return f"Rp{val / 1e6:,.0f}Jt"

    x_hdr = "SALES GROWTH" if x_axis_metric == "Qty Growth (%)" else "GROSS SALES (ABSOLUTE)"

    html_content = textwrap.dedent(f"""
        <style>
        .grid-container {{ display: grid; grid-template-columns: 100px 1fr 1fr 1fr; grid-template-rows: 25px 25px 1fr 1fr 1fr; gap: 12px; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin-bottom: 2rem; margin-top: 1rem; }}
        .box {{ padding: 15px; color: white; border-radius: 2px; display: flex; flex-direction: column; justify-content: space-between; min-height: 145px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }}
        .box-title {{ font-size: 13px; font-weight: 600; opacity: 0.95; }}
        .box-main {{ font-size: 32px; font-weight: 700; margin: 8px 0; }}
        .box-sub {{ font-size: 13px; font-weight: 400; opacity: 0.9; }}
        .box-sub-small {{ font-size: 12px; font-weight: 500; opacity: 0.95; line-height: 1.5; }}
        .bg-blue {{ background-color: #3871b6; }}
        .bg-green {{ background-color: #319b5e; }}
        .bg-red {{ background-color: #c03d32; }}
        .bg-gold {{ background-color: #d89f0e; }}
        .hdr-main {{ grid-column: 1 / 5; text-align: center; font-weight: 700; font-size: 13px; color: #4b5563; letter-spacing: 1px; }}
        .hdr-col {{ text-align: center; color: #6b7280; font-size: 13px; font-weight: 700; align-self: end; }}
        .hdr-row {{ text-align: center; font-size: 12px; font-weight: 700; color: #6b7280; display: flex; flex-direction: column; justify-content: center; }}
        .axis-label-y {{ writing-mode: vertical-rl; transform: rotate(180deg); position: absolute; left: 0; font-size: 12px; letter-spacing: 1px; color: #4b5563; font-weight: 700; height: 100%; text-align: center; }}
        .legend-container {{ display: flex; justify-content: center; gap: 40px; margin-top: 10px; font-size: 13px; font-weight: 700; color: #374151; }}
        .legend-item {{ display: flex; align-items: center; gap: 8px; }}
        .legend-box {{ width: 25px; height: 15px; border-radius: 2px; }}
        </style>

        <div style="position: relative;">
        <div class="axis-label-y">{margin_type.upper()}</div>
        <div class="grid-container" style="margin-left: 30px;">
        <div class="hdr-main">{x_hdr} &rarr;</div>
        <div></div>
        <div class="hdr-col">Low<br>(&lt;{format_x_val(x_low_thresh)})</div>
        <div class="hdr-col">Med<br>({format_x_val(x_low_thresh)} &ndash; {format_x_val(x_high_thresh)})</div>
        <div class="hdr-col">High<br>(&gt;{format_x_val(x_high_thresh)})</div>

        <div class="hdr-row">High<br>&gt;{y_high_thresh:.1f}%</div>
        {get_box_html('Box 1', 'bg-blue')}
        {get_box_html('Box 4', 'bg-blue')}
        {get_box_html('Box 7', 'bg-green')}

        <div class="hdr-row">Med<br>{y_low_thresh:.1f}&ndash;{y_high_thresh:.1f}%</div>
        {get_box_html('Box 2', 'bg-gold')}
        {get_box_html('Box 5', 'bg-gold')}
        {get_box_html('Box 8', 'bg-green')}

        <div class="hdr-row">Low<br>&lt;{y_low_thresh:.1f}%</div>
        {get_box_html('Box 3', 'bg-red')}
        {get_box_html('Box 6', 'bg-gold')}
        {get_box_html('Box 9', 'bg-gold')}
        </div>
        </div>

        <div class="legend-container">
        <div class="legend-item"><div class="legend-box bg-green"></div> Grow</div>
        <div class="legend-item"><div class="legend-box bg-blue"></div> Keep</div>
        <div class="legend-item"><div class="legend-box bg-gold"></div> Fix</div>
        <div class="legend-item"><div class="legend-box bg-red"></div> Exit candidate</div>
        </div>
        <br>
    """)
    st.markdown(html_content, unsafe_allow_html=True)
    return html_content


# ==========================================
# PPTX GENERATION ENGINE
# ==========================================
def create_portfolio_presentation(df_main, margin_type, margin_val_col):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(22, 54, 92)
    GRAY_TEXT = RGBColor(89, 89, 89)
    ACCENT_BLUE = RGBColor(0, 112, 192)
    LIGHT_BG = RGBColor(242, 242, 242)

    def add_consulting_slide(title_text, subtitle_text, metrics, insights):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_TEXT

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = NAVY
        line.line.color.rgb = NAVY

        box_width = 4.0
        for i, metric in enumerate(metrics):
            left_pos = 0.5 + (i * (box_width + 0.165))
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.3),
                                           Inches(box_width), Inches(0.6))
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_BG
            shape.line.color.rgb = ACCENT_BLUE

            tf_m = shape.text_frame
            tf_m.vertical_anchor = 3
            p_m = tf_m.paragraphs[0]
            p_m.alignment = PP_ALIGN.CENTER
            p_m.text = metric
            p_m.font.bold = True
            p_m.font.size = Pt(14)
            p_m.font.color.rgb = NAVY

        chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(7.5), Inches(4.8))
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = LIGHT_BG
        chart_box.line.color.rgb = GRAY_TEXT

        tf_c = chart_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = f"[ Insert {margin_type} Matrix Visualization Here ]\n(Copy-paste directly from the dashboard)"
        p_c.font.size = Pt(14)
        p_c.font.color.rgb = GRAY_TEXT
        p_c.font.italic = True

        insight_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.6), Inches(5.0))
        tf_insight = insight_box.text_frame
        tf_insight.word_wrap = True

        p_ins_head = tf_insight.paragraphs[0]
        p_ins_head.text = "Executive Insights & Action Mandate"
        p_ins_head.font.bold = True
        p_ins_head.font.size = Pt(16)
        p_ins_head.font.color.rgb = NAVY

        for insight in insights:
            parts = insight.split(":", 1)
            p_i = tf_insight.add_paragraph()
            p_i.space_before = Pt(12)
            p_i.level = 0

            if len(parts) == 2:
                run_bold = p_i.add_run()
                run_bold.text = f"• {parts[0]}:"
                run_bold.font.bold = True
                run_bold.font.size = Pt(12)
                run_bold.font.color.rgb = NAVY

                run_text = p_i.add_run()
                run_text.text = parts[1]
                run_text.font.size = Pt(12)
                run_text.font.color.rgb = GRAY_TEXT
            else:
                run_text = p_i.add_run()
                run_text.text = f"• {insight}"
                run_text.font.size = Pt(12)
                run_text.font.color.rgb = GRAY_TEXT

        return slide

    total_rev = df_main['Gross Sales (Current)'].sum()
    total_margin = df_main[margin_val_col].sum()
    blended_margin = (total_margin / total_rev * 100) if total_rev > 0 else 0
    total_skus = len(df_main)

    rev_str = f"Rp {total_rev / 1e12:.2f} Trillion" if total_rev >= 1e12 else f"Rp {total_rev / 1e9:.2f} Billion"

    add_consulting_slide(
        title_text=f"Global Portfolio Rationalization ({margin_type} View)",
        subtitle_text="Comprehensive assessment highlighting significant margin dilution and strategic resource allocation.",
        metrics=[
            f"Total Active SKUs: {total_skus:,}",
            f"Total Revenue: {rev_str}",
            f"Blended {margin_type}: {blended_margin:.1f}%"
        ],
        insights=[
            "Value-Destructive Tail: A significant portion of SKUs exhibit severe margin compression. Immediate rationalization or price restructuring is mandated.",
            "Aggressive Pruning & Pricing Gate: Execute targeted price hikes for underperformers. SKUs failing to clear the profitability hurdle post-hike must be slated for phase-out.",
            "Core Growth Engine: Flagship SKUs continue to deliver superior margins. Capital allocation should immediately pivot toward expanding distribution for these core assets."
        ]
    )

    return prs


# ==========================================
# MAIN APPLICATION LOGIC
# ==========================================
def main():
    st.markdown("""
        <div style="background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%); 
                    padding: 1.2rem; border-radius: 15px; margin-bottom: 1rem; color: white;
                    text-align: center;">
            <h1 style="margin:0; color: #38bdf8; font-size: 1.8rem; font-weight: 800;">
                📊 UNIFIED 9-BOX PORTFOLIO MATRIX
            </h1>
            <p style="margin: 0.4rem 0 0 0; opacity: 0.9; font-size: 1rem;">
                Strategic Engine for Gross Margin & Contribution Margin Optimization
            </p>
        </div>
    """, unsafe_allow_html=True)

    # ---------------------------------------------------------
    # FILE UPLOADER
    # ---------------------------------------------------------
    st.markdown("### 📂 Upload Data Sumber")
    st.info(
        "Silakan unggah file Excel **Data 9-box** (misal: Data 9-box CM.xlsx) dari komputer Anda untuk memulai analisis.")

    uploaded_file = st.file_uploader("Unggah File Data 9-Box (.xlsx)", type=["xlsx"])

    if uploaded_file is None:
        st.stop()

    # Load and preserve the raw data universally
    df_raw = load_9box_data(uploaded_file)

    if df_raw.empty:
        st.warning(
            "⚠️ File tidak dapat dibaca atau formatnya salah. Pastikan file memiliki sheet 'LOCAL' dan 'EXPORT'.")
        st.stop()
        
    df_main = df_raw.copy()

    # ---------------------------------------------------------
    # SKU CONSOLIDATION VIEW TOGGLE 
    # ---------------------------------------------------------
    st.markdown("### 🔀 SKU Consolidation View")
    view_mode = st.radio(
        "Select Data Granularity:",
        [
            "Uncommon (Original Granular SKUs)", 
            "Common (Consolidated Master SKUs)",
            "Commonized Only (Master C- Prefix SKUs)",
            "Uncommonized Only (Granular C- Prefix SKUs)"
        ],
        horizontal=True,
        help="Switch to 'Common' to group data by Master SKU. Choose 'Commonized Only' for merged items, or 'Uncommonized Only' to see the original granular details of those merged items."
    )

    # Grouping Logic for "Common" and "Commonized Only" views
    if view_mode in ["Common (Consolidated Master SKUs)", "Commonized Only (Master C- Prefix SKUs)"]:
        if 'New Code' in df_main.columns and 'New Product Name' in df_main.columns:
            df_main['New Code'] = df_main['New Code'].astype(str)
            df_main['New Product Name'] = df_main['New Product Name'].astype(str)

            # Columns required for grouping
            group_cols = ['Source_Sheet', 'New Code', 'New Product Name']
            
            # Identify numeric columns for summation
            num_cols = df_main.select_dtypes(include=[np.number]).columns.tolist()
            sum_cols = [c for c in num_cols if not c.endswith('(%)')]
            
            agg_dict = {c: 'sum' for c in sum_cols}
            
            # Carry over categorical columns
            cat_cols = ['Remark', 'Status', 'Country']
            for cat in cat_cols:
                if cat in df_main.columns:
                    agg_dict[cat] = lambda x: next(iter(x.dropna()), '')

            # Execute dynamic GroupBy
            df_main = df_main.groupby(group_cols, as_index=False).agg(agg_dict)

            # Re-calculate Percentages based on summed values
            if 'Qty (Current)' in df_main.columns and 'Qty (Previous)' in df_main.columns:
                df_main['Qty Growth (%)'] = np.where(
                    df_main['Qty (Previous)'] != 0,
                    ((df_main['Qty (Current)'] - df_main['Qty (Previous)']) / df_main['Qty (Previous)'].abs()) * 100,
                    0.0
                )
            
            if 'Gross Sales (Current)' in df_main.columns:
                sales = df_main['Gross Sales (Current)']
                if 'Gross Profit (Current)' in df_main.columns:
                    df_main['Gross Profit (%)'] = np.where(sales > 0, (df_main['Gross Profit (Current)'] / sales) * 100, np.where(df_main['Gross Profit (Current)'] < 0, -100.0, 0.0))
                if 'Gross Margin (Current)' in df_main.columns:
                    df_main['Gross Margin (%)'] = np.where(sales > 0, (df_main['Gross Margin (Current)'] / sales) * 100, np.where(df_main['Gross Margin (Current)'] < 0, -100.0, 0.0))
                if 'Contribution Margin (Current)' in df_main.columns:
                    df_main['Contribution Margin (%)'] = np.where(sales > 0, (df_main['Contribution Margin (Current)'] / sales) * 100, np.where(df_main['Contribution Margin (Current)'] < 0, -100.0, 0.0))

            # Replace Product Name with Master Name
            df_main['Product Name'] = df_main['New Product Name']
            
            # --- COMMONIZED ONLY FILTER ---
            if view_mode == "Commonized Only (Master C- Prefix SKUs)":
                df_main = df_main[df_main['Product Name'].str.startswith('C-', na=False)]
            
        else:
            st.warning("⚠️ Columns 'New Code' and 'New Product Name' are missing in the uploaded dataset.")

    # Filtering Logic for "Uncommonized Only" view (Keeps data Granular)
    elif view_mode == "Uncommonized Only (Granular C- Prefix SKUs)":
        if 'New Product Name' in df_main.columns:
            # Hanya tampilkan data granular yang New Product Name-nya berawalan "C-"
            df_main = df_main[df_main['New Product Name'].astype(str).str.startswith('C-')]
        else:
            st.warning("⚠️ Column 'New Product Name' is missing in the uploaded dataset.")

    st.markdown("---")

    tab_main_matrix, tab_b3_intersection, tab_progress = st.tabs([
        "📈 MAIN 9-BOX ANALYSIS", 
        "⚠️ BOX 3 DEEP DIVE", 
        "📉 RATIONALIZATION PROGRESS"
    ])

    # ---------------------------------------------------------
    # TAB 1: MAIN 9-BOX ANALYSIS
    # ---------------------------------------------------------
    with tab_main_matrix:
        st.markdown("### 🔍 Filter & Metric Configuration")
        
        col1, col2, col3 = st.columns(3)

        with col1:
            market_filter = st.selectbox(
                "Select Market Focus:",
                ["ALL (Local + Export)", "LOCAL", "EXPORT"]
            )

        with col2:
            margin_selector = st.selectbox(
                "Y-Axis Metric (Margin):",
                ["Gross Profit", "Gross Margin", "Contribution Margin"]
            )

        with col3:
            x_axis_selector = st.selectbox(
                "X-Axis Metric (Performance):",
                ["Gross Sales (Current)", "Qty Growth (%)"],
                help="Pilih 'Gross Sales (Current)' untuk nominal absolut, atau 'Qty Growth (%)' untuk melihat pertumbuhan."
            )

        col4, col5, col6, col7 = st.columns(4)
        
        with col4:
            bubble_size_selector = st.selectbox(
                "Bubble Size Metric:",
                ["Gross Sales (Current)", "Gross Profit (Current)", "Gross Margin (Current)",
                 "Contribution Margin (Current)", "Qty (Current)"]
            )

        with col5:
            if 'Remark' in df_main.columns:
                valid_remarks = sorted(
                    [str(x) for x in df_main['Remark'].unique() if pd.notna(x) and str(x).strip() != ''])
                remark_options = ["ALL"] + valid_remarks
            else:
                remark_options = ["ALL"]

            remark_filter = st.multiselect("Select Remark(s):", remark_options, default=["ALL"])

        with col6:
            if 'Status' in df_main.columns:
                valid_statuses = sorted(
                    [str(x) for x in df_main['Status'].unique() if pd.notna(x) and str(x).strip() != ''])
                status_options = ["ALL"] + valid_statuses
            else:
                status_options = ["ALL"]
                
            status_filter = st.multiselect("Select Status(es):", status_options, default=["ALL"])

        with col7:
            sku_search = st.text_input("Specific SKU Search:", "")

        if margin_selector == "Gross Profit":
            y_col = 'Gross Profit (%)'
            margin_val_col = 'Gross Profit (Current)'
        elif margin_selector == "Gross Margin":
            y_col = 'Gross Margin (%)'
            margin_val_col = 'Gross Margin (Current)'
        else:
            y_col = 'Contribution Margin (%)'
            margin_val_col = 'Contribution Margin (Current)'

        x_col = x_axis_selector

        filtered_df = df_main.copy()
        
        if market_filter != "ALL (Local + Export)":
            filtered_df = filtered_df[filtered_df['Source_Sheet'] == market_filter]

        if remark_filter and 'Remark' in filtered_df.columns:
            actual_remarks = [r for r in remark_filter if r != "ALL"]
            if actual_remarks:
                filtered_df = filtered_df[filtered_df['Remark'].astype(str).isin(actual_remarks)]

        if status_filter and 'Status' in filtered_df.columns:
            actual_statuses = [s for s in status_filter if s != "ALL"]
            if actual_statuses:
                filtered_df = filtered_df[filtered_df['Status'].astype(str).isin(actual_statuses)]

        if sku_search:
            search_col = 'Product Name' if 'Product Name' in filtered_df.columns else (
                'SKU' if 'SKU' in filtered_df.columns else None)
            if search_col:
                filtered_df = filtered_df[
                    filtered_df[search_col].astype(str).str.contains(sku_search, case=False, na=False)]

        st.markdown("### 🎚️ Smart-Scaling & Outlier Control")

        col_out1, col_out2 = st.columns(2)
        with col_out1:
            min_outlier_limit_x = st.number_input(
                f"Batas MINIMAL {x_axis_selector} untuk Rata-rata (Scope-out Long-tail):",
                value=0.0,
                step=1000000.0 if x_axis_selector == "Gross Sales (Current)" else 1.0,
                help="SKU di Bawah batas ini TETAP TAMPIL, tapi TIDAK DIAJAK menghitung nilai Rata-rata (Center) Sumbu X."
            )

        with col_out2:
            scale_factor_x = st.number_input(
                f"Skala Visual Box Kanan (Multiplier X-Axis):",
                value=0.2,
                step=0.1,
                min_value=0.001,
                max_value=10.0,
                help="Set ke 1.0 untuk melihat grafik asli (bisa gepeng jika ada outlier). Set < 1.0 (misal 0.2 atau 0.1) untuk mengecilkan visual box kanan (B7,8,9) agar box kiri dan tengah mendapat porsi ruang layar lebih luas tanpa menyembunyikan SKU apa pun."
            )

        if not filtered_df.empty:
            positive_df = filtered_df[filtered_df[margin_val_col] >= 0]
            total_sales_pos = positive_df['Gross Sales (Current)'].sum()

            if total_sales_pos > 0:
                avg_margin_pct = (positive_df[margin_val_col].sum() / total_sales_pos) * 100
            else:
                avg_margin_pct = positive_df[y_col].mean() if not positive_df.empty else 25.0

            if pd.isna(avg_margin_pct) or avg_margin_pct <= 0:
                avg_margin_pct = 1.0

            avg_y = avg_margin_pct
            def_y_low = avg_y * (2.0 / 3.0)
            def_y_high = avg_y * (4.0 / 3.0)

            normal_x_df = filtered_df[filtered_df[x_col] >= min_outlier_limit_x]
            if not normal_x_df.empty:
                robust_avg_x = normal_x_df[x_col].mean()
            else:
                robust_avg_x = filtered_df[x_col].mean()

            if pd.isna(robust_avg_x) or robust_avg_x <= 0:
                robust_avg_x = 100.0

            def_x_low = robust_avg_x * (2.0 / 3.0)
            def_x_high = robust_avg_x * (4.0 / 3.0)

            x_max_raw = filtered_df[x_col].max()
            x_min_raw = filtered_df[x_col].min()
            x_span_raw = x_max_raw - x_min_raw if x_max_raw != x_min_raw else robust_avg_x

            plot_x_min = x_min_raw - (abs(x_span_raw) * 0.05)
            if x_min_raw < 0:
                plot_x_min = x_min_raw - (abs(x_span_raw) * 0.05)
            plot_x_max = x_max_raw + (abs(x_span_raw) * 0.05)

            y_min_raw = filtered_df[y_col].min()
            y_max_raw = filtered_df[y_col].max()
            y_span_raw = y_max_raw - y_min_raw if y_max_raw != y_min_raw else 100.0

            plot_y_min = y_min_raw - (y_span_raw * 0.05)
            plot_y_max = max(y_max_raw, def_y_high * 1.5) + (y_span_raw * 0.05)
            plot_y_min = min(plot_y_min, def_y_low - (def_y_low * 0.5))
            plot_y_max = max(plot_y_max, def_y_high + (def_y_high * 0.5))
        else:
            def_y_low, def_y_high = 16.67, 33.33
            def_x_low, def_x_high = 66.67, 133.33
            robust_avg_x, avg_margin_pct = 100.0, 25.0
            plot_x_min, plot_x_max, plot_y_min, plot_y_max = -5, 200, -20, 50

        with st.form("threshold_form"):
            st.info(
                "💡 **INFO:** Sumbu X dan Y menggunakan nilai Rata-Rata sebagai titik tengah mutlak. Garis Rata-rata ditandai dengan warna Merah. Garis Threshold Atas dan Bawah otomatis disesuaikan secara simetris terhadap Rata-rata.")

            # HASH KEY UNTUK MENGAKALI CACHE FORM INPUT
            filter_state_str = f"{view_mode}_{market_filter}_{margin_selector}_{x_axis_selector}_{remark_filter}_{status_filter}_{sku_search}_{min_outlier_limit_x}"
            form_key_suffix = hashlib.md5(filter_state_str.encode('utf-8')).hexdigest()

            step_x_input = float(abs(def_x_high - def_x_low) / 2) if def_x_high != def_x_low else 1.0

            col_tx1, col_tx2, col_ty1, col_ty2 = st.columns(4)
            with col_tx1:
                x_low_thresh_input = st.number_input(f"X-Axis Low to Med", value=float(def_x_low), step=step_x_input, key=f"xl_{form_key_suffix}")
            with col_tx2:
                x_high_thresh_input = st.number_input(f"X-Axis Med to High", value=float(def_x_high), step=step_x_input, key=f"xh_{form_key_suffix}")
            with col_ty1:
                y_low_thresh_input = st.number_input("Y-Axis Low to Med (%)", value=float(def_y_low), step=1.0, key=f"yl_{form_key_suffix}")
            with col_ty2:
                y_high_thresh_input = st.number_input("Y-Axis Med to High (%)", value=float(def_y_high), step=1.0, key=f"yh_{form_key_suffix}")

            run_thresholds = st.form_submit_button("▶ RUN & UPDATE MATRIX", type="primary")

        y_low_thresh = float(y_low_thresh_input)
        y_high_thresh = float(y_high_thresh_input)
        x_low_thresh = float(x_low_thresh_input)
        x_high_thresh = float(x_high_thresh_input)

        if filtered_df.empty:
            st.error("No data found matching the selected filters.")
        else:
            st.markdown("---")
            st.subheader(f"📈 Matrix View: {market_filter} | {margin_selector} vs {x_axis_selector}")

            x_lbl_low = "Low Growth" if x_axis_selector == "Qty Growth (%)" else "Low Sales"
            x_lbl_med = "Med Growth" if x_axis_selector == "Qty Growth (%)" else "Med Sales"
            x_lbl_high = "High Growth" if x_axis_selector == "Qty Growth (%)" else "High Sales"

            b1_id = f'Box 1 (High Margin, {x_lbl_low})'
            b2_id = f'Box 2 (Med Margin, {x_lbl_low})'
            b3_id = f'Box 3 (Low Margin, {x_lbl_low})'
            b4_id = f'Box 4 (High Margin, {x_lbl_med})'
            b5_id = f'Box 5 (Med Margin, {x_lbl_med})'
            b6_id = f'Box 6 (Low Margin, {x_lbl_med})'
            b7_id = f'Box 7 (High Margin, {x_lbl_high})'
            b8_id = f'Box 8 (Med Margin, {x_lbl_high})'
            b9_id = f'Box 9 (Low Margin, {x_lbl_high})'

            conditions_9box = [
                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] < x_low_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] < x_low_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] < x_low_thresh),

                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] >= x_low_thresh) & (
                        filtered_df[x_col] <= x_high_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] >= x_low_thresh) & (filtered_df[x_col] <= x_high_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] >= x_low_thresh) & (
                        filtered_df[x_col] <= x_high_thresh),

                (filtered_df[y_col] > y_high_thresh) & (filtered_df[x_col] > x_high_thresh),
                (filtered_df[y_col] >= y_low_thresh) & (filtered_df[y_col] <= y_high_thresh) & (
                        filtered_df[x_col] > x_high_thresh),
                (filtered_df[y_col] < y_low_thresh) & (filtered_df[x_col] > x_high_thresh)
            ]

            choices_9box = [b1_id, b2_id, b3_id, b4_id, b5_id, b6_id, b7_id, b8_id, b9_id]

            filtered_df['Dynamic 9-Box Category'] = np.select(
                conditions_9box, choices_9box, default=b6_id
            )

            render_9box_summary_grid(filtered_df, margin_selector, margin_val_col, y_low_thresh, y_high_thresh,
                                     x_axis_selector, x_low_thresh, x_high_thresh)

            # ==========================================
            # SYNCED BUBBLE CHART RENDERING
            # ==========================================
            plot_df = filtered_df.copy()

            if bubble_size_selector in plot_df.columns:
                plot_df['Bubble_Size'] = pd.to_numeric(plot_df[bubble_size_selector], errors='coerce').fillna(
                    0).abs().replace(0, 1)
            else:
                plot_df['Bubble_Size'] = 10

            box_counts = plot_df['Dynamic 9-Box Category'].value_counts().to_dict()
            total_items = len(plot_df)

            st.info(f"💡 **Total Evaluated & Plotted**: {total_items} SKUs dynamically distributed across the Matrix.")

            name_col = 'Product Name' if 'Product Name' in plot_df.columns else (
                'SKU' if 'SKU' in plot_df.columns else plot_df.columns[0])

            # Color Configuration (Adding Black for Commonized SKUs)
            color_discrete_map = {
                b1_id: '#3871b6', b4_id: '#3871b6', b7_id: '#319b5e',
                b8_id: '#319b5e', b2_id: '#d89f0e', b5_id: '#d89f0e',
                b6_id: '#d89f0e', b9_id: '#d89f0e', b3_id: '#c03d32',
                'Commonized SKU (Black)': '#000000'
            }

            # Map the Box colors first, then override Commonized SKUs to Black
            plot_df['Plot_Color_Category'] = plot_df['Dynamic 9-Box Category']
            if 'New Product Name' in plot_df.columns:
                mask_commonized = plot_df['New Product Name'].fillna('').astype(str).str.startswith('C-')
                plot_df.loc[mask_commonized, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

            def format_x_val_str(val):
                if x_axis_selector == "Qty Growth (%)": return f"{val:.1f}%"
                if val >= 1e9 or val <= -1e9: return f"Rp{val / 1e9:,.1f}M"
                return f"Rp{val / 1e6:,.0f}Jt"

            x_format = ':.2f' if x_axis_selector == 'Qty Growth (%)' else ':,.0f'

            def apply_custom_x_scale(v):
                if v <= x_high_thresh:
                    return v
                else:
                    return x_high_thresh + (v - x_high_thresh) * scale_factor_x

            plot_df['X_Plot'] = plot_df[x_col].apply(apply_custom_x_scale)
            plot_df['Y_Plot'] = plot_df[y_col]

            hover_data_dict = {
                'Source_Sheet': True if 'Source_Sheet' in plot_df.columns else False,
                'Dynamic 9-Box Category': True, 
                'Plot_Color_Category': False,
                y_col: ':.2f',
                x_col: x_format,
                'X_Plot': False,
                'Y_Plot': False,
                'Bubble_Size': False,
                bubble_size_selector: ':,.0f'
            }

            if 'Qty (Current)' in plot_df.columns: hover_data_dict['Qty (Current)'] = ':,.0f'

            fig = px.scatter(
                plot_df,
                x='X_Plot',
                y='Y_Plot',
                size="Bubble_Size",
                color="Plot_Color_Category",
                color_discrete_map=color_discrete_map,
                hover_name=name_col,
                hover_data=hover_data_dict,
                title=f"9-Box Bubble Chart: {margin_selector} vs {x_axis_selector} (Total: {total_items} SKUs)",
                size_max=60,
                render_mode="svg"
            )

            fig.add_vline(x=apply_custom_x_scale(x_low_thresh), line_dash="dash", line_color="black", line_width=2,
                          opacity=0.8, annotation_text=f" Low ({format_x_val_str(x_low_thresh)})",
                          annotation_position="top left")
            fig.add_vline(x=apply_custom_x_scale(x_high_thresh), line_dash="dash", line_color="black", line_width=2,
                          opacity=0.8, annotation_text=f" High ({format_x_val_str(x_high_thresh)})",
                          annotation_position="top right")
            fig.add_hline(y=y_low_thresh, line_dash="dash", line_color="black", line_width=2, opacity=0.8,
                          annotation_text=f"Low ({y_low_thresh:.1f}%)", annotation_position="bottom right")
            fig.add_hline(y=y_high_thresh, line_dash="dash", line_color="black", line_width=2, opacity=0.8,
                          annotation_text=f"High ({y_high_thresh:.1f}%)", annotation_position="top right")

            fig.add_vline(x=apply_custom_x_scale(robust_avg_x), line_dash="dot", line_color="red", line_width=2,
                          opacity=0.6, annotation_text=f" ← AVG X ({format_x_val_str(robust_avg_x)})",
                          annotation_position="top right", annotation_font_color="red")
            fig.add_hline(y=avg_margin_pct, line_dash="dot", line_color="red", line_width=2, opacity=0.6,
                          annotation_text=f"AVG Y ({avg_margin_pct:.1f}%)", annotation_position="top right",
                          annotation_font_color="red")

            mapped_plot_x_min = apply_custom_x_scale(plot_x_min)
            mapped_plot_x_max = apply_custom_x_scale(plot_x_max)
            mapped_x_low = apply_custom_x_scale(x_low_thresh)
            mapped_x_high = apply_custom_x_scale(x_high_thresh)

            mid_x_low = (mapped_plot_x_min + mapped_x_low) / 2
            mid_x_med = (mapped_x_low + mapped_x_high) / 2
            mid_x_high = (mapped_x_high + mapped_plot_x_max) / 2

            mid_y_low = (plot_y_min + y_low_thresh) / 2
            mid_y_med = (y_low_thresh + y_high_thresh) / 2
            mid_y_high = (y_high_thresh + plot_y_max) / 2

            box_coords = {
                b1_id: {'x': mid_x_low, 'y': mid_y_high, 'id': 'B1'},
                b2_id: {'x': mid_x_low, 'y': mid_y_med, 'id': 'B2'},
                b3_id: {'x': mid_x_low, 'y': mid_y_low, 'id': 'B3'},
                b4_id: {'x': mid_x_med, 'y': mid_y_high, 'id': 'B4'},
                b5_id: {'x': mid_x_med, 'y': mid_y_med, 'id': 'B5'},
                b6_id: {'x': mid_x_med, 'y': mid_y_low, 'id': 'B6'},
                b7_id: {'x': mid_x_high, 'y': mid_y_high, 'id': 'B7'},
                b8_id: {'x': mid_x_high, 'y': mid_y_med, 'id': 'B8'},
                b9_id: {'x': mid_x_high, 'y': mid_y_low, 'id': 'B9'}
            }

            for box_name, coords in box_coords.items():
                count = box_counts.get(box_name, 0)
                if count > 0:
                    is_b5 = (coords['id'] == 'B5')
                    fig.add_annotation(
                        x=coords['x'],
                        y=coords['y'],
                        ax=mapped_x_high + (mapped_plot_x_max - mapped_x_high) * 0.05 if is_b5 else None,
                        ay=mid_y_low if is_b5 else None,
                        axref="x" if is_b5 else None,
                        ayref="y" if is_b5 else None,
                        xref="x",
                        yref="y",
                        text=f"<span style='color:#1f2937;'><b>{coords['id']}</b></span><br><b>{count} SKUs</b>",
                        showarrow=is_b5,
                        arrowhead=2 if is_b5 else 0,
                        arrowsize=1 if is_b5 else 1,
                        arrowwidth=1 if is_b5 else 0.1,
                        arrowcolor="#1f2937" if is_b5 else None,
                        standoff=0,
                        xanchor="left" if is_b5 else "center",
                        font=dict(size=13, color="#374151"),
                        bgcolor="rgba(255, 255, 255, 0.9)" if is_b5 else "rgba(255, 255, 255, 0.8)",
                        bordercolor="rgba(15, 23, 42, 0.8)" if is_b5 else "rgba(15, 23, 42, 0.3)",
                        borderwidth=2 if is_b5 else 1,
                        borderpad=4
                    )

            fig.add_annotation(
                x=0.98,
                y=0.98,
                xref="paper",
                yref="paper",
                text=f"<b>TOTAL:<br>{total_items} SKUs</b>",
                showarrow=False,
                font=dict(size=13, color="white"),
                bgcolor="#374151",
                bordercolor="black",
                borderwidth=1,
                borderpad=5
            )

            ticks_x_actual = [plot_x_min, 0, x_low_thresh, robust_avg_x, x_high_thresh, plot_x_max]
            ticks_x_actual = sorted(list(set([v for v in ticks_x_actual if v >= plot_x_min and v <= plot_x_max])))

            tickvals_x = [apply_custom_x_scale(v) for v in ticks_x_actual]
            ticktext_x = [format_x_val_str(v) for v in ticks_x_actual]

            fig.update_xaxes(
                range=[mapped_plot_x_min, mapped_plot_x_max],
                tickvals=tickvals_x,
                ticktext=ticktext_x,
                title_text="LTM Quantity Growth Rate (%)" if x_axis_selector == "Qty Growth (%)" else "Gross Sales Absolute (IDR)"
            )
            fig.update_yaxes(range=[plot_y_min, plot_y_max], title_text=f"{margin_selector} (%)")

            fig.update_layout(
                height=550,
                hovermode="closest",
                showlegend=True,
                legend=dict(
                    title=None,
                    orientation="h",
                    yanchor="top",
                    y=-0.15,
                    xanchor="center",
                    x=0.5
                ),
                margin=dict(t=50, b=50, l=50, r=50)
            )

            st.plotly_chart(fig, use_container_width=True, config={
                'displayModeBar': True,
                'modeBarButtonsToAdd': ['toImage'],
                'displaylogo': False,
                'toImageButtonOptions': {
                    'format': 'png',
                    'filename': f'9Box_Matrix_{margin_selector.replace(" ", "_")}',
                    'height': 800,
                    'width': 1400,
                    'scale': 2
                }
            })

            # ---------------------------------------------------------
            # NEW: 9-BOX INTERACTIVE QUICK EXPORT GRID
            # ---------------------------------------------------------
            st.markdown("---")
            st.subheader("📥 Interactive 9-Box Quick Export")
            st.info("💡 **Quick Action:** Klik tombol di bawah ini untuk langsung men-download isi dari masing-masing Box. Angka finansial diekspor murni sebagai *Raw Numbers* (agar bisa di-SUM di Excel), sementara kode produk dan teks lainnya dijamin aman sebagai String.")
            
            def get_raw_excel(df_in):
                # Buang kolom teknis dari Plotly agar rapi persis seperti tabel "Data Detail & Export"
                df_out = df_in.drop(columns=['Bubble_Size', 'X_Plot', 'Y_Plot', 'Plot_Color_Category', 'Dynamic 9-Box Category'], errors='ignore').copy()
                
                # Pastikan kolom bertipe teks tetap dibaca string oleh Excel (mencegah auto-convert kode produk jadi angka)
                text_cols = ['SKU', 'Product Name', 'New Code', 'New Product Name', 'Remark', 'Status', 'Country', 'Source_Sheet']
                for col in text_cols:
                    if col in df_out.columns:
                        df_out[col] = df_out[col].astype(str)
                
                buffer = io.BytesIO()
                with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
                    df_out.to_excel(writer, index=False, sheet_name='Box_Export')
                    
                    # Auto-adjust column width for neatness di Excel
                    worksheet = writer.sheets['Box_Export']
                    for col in worksheet.columns:
                        max_length = 0
                        column = col[0].column_letter
                        for cell in col:
                            try:
                                if len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        adjusted_width = (max_length + 2)
                        worksheet.column_dimensions[column].width = adjusted_width
                        
                return buffer.getvalue()

            grid_col1, grid_col2, grid_col3 = st.columns(3)
            grid_col4, grid_col5, grid_col6 = st.columns(3)
            grid_col7, grid_col8, grid_col9 = st.columns(3)

            boxes_layout = [
                (grid_col1, b1_id, "Box 1"), (grid_col2, b4_id, "Box 4"), (grid_col3, b7_id, "Box 7"),
                (grid_col4, b2_id, "Box 2"), (grid_col5, b5_id, "Box 5"), (grid_col6, b8_id, "Box 8"),
                (grid_col7, b3_id, "Box 3"), (grid_col8, b6_id, "Box 6"), (grid_col9, b9_id, "Box 9")
            ]

            for col, box_id, label in boxes_layout:
                box_df = filtered_df[filtered_df['Dynamic 9-Box Category'] == box_id].copy()
                count = len(box_df)
                with col:
                    if count > 0:
                        excel_data = get_raw_excel(box_df)
                        st.download_button(
                            label=f"📥 Download {label} ({count} SKUs)",
                            data=excel_data,
                            file_name=f"Export_{label.replace(' ', '')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                    else:
                        st.button(f"📥 {label} (0 SKUs)", disabled=True, use_container_width=True)


            # ---------------------------------------------------------
            # DATA DETAIL TABLE & EXPORT
            # ---------------------------------------------------------
            st.markdown("---")
            col_hdr1, col_hdr2, col_hdr3, col_hdr4 = st.columns([1.5, 1, 1, 1])
            with col_hdr1:
                st.subheader("📋 Data Detail & Export")

            display_df = filtered_df.drop(columns=['Bubble_Size', 'X_Plot', 'Y_Plot', 'Plot_Color_Category'], errors='ignore')

            buffer_excel_main = io.BytesIO()
            with pd.ExcelWriter(buffer_excel_main, engine='openpyxl') as writer:
                display_df.to_excel(writer, index=False, sheet_name='9Box_Data')

            with col_hdr2:
                st.download_button(
                    label="📥 Download Excel",
                    data=buffer_excel_main.getvalue(),
                    file_name=f"9Box_Matrix_{margin_selector.replace(' ', '_')}_{market_filter}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    type="primary",
                    use_container_width=True
                )

            with col_hdr3:
                prs = create_portfolio_presentation(filtered_df, margin_selector, margin_val_col)
                buffer_ppt = io.BytesIO()
                prs.save(buffer_ppt)

                st.download_button(
                    label="📥 Download Deck (.pptx)",
                    data=buffer_ppt.getvalue(),
                    file_name=f"Consulting_Deck_{margin_selector.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="secondary",
                    use_container_width=True
                )

            with col_hdr4:
                buffer_html = io.StringIO()
                fig.write_html(buffer_html, include_plotlyjs='cdn')
                st.download_button(
                    label="📥 Download Chart (.html)",
                    data=buffer_html.getvalue(),
                    file_name=f"9Box_Interactive_Chart_{margin_selector.replace(' ', '_')}.html",
                    mime="text/html",
                    type="secondary",
                    use_container_width=True
                )

            format_dict = {}
            if 'Gross Sales (Current)' in display_df.columns: format_dict['Gross Sales (Current)'] = '{:,.0f}'
            if 'Return (Current)' in display_df.columns: format_dict['Return (Current)'] = '{:,.0f}'
            if 'Qty (Current)' in display_df.columns: format_dict['Qty (Current)'] = '{:,.0f}'
            if 'Gross Sales (Previous)' in display_df.columns: format_dict['Gross Sales (Previous)'] = '{:,.0f}'
            if 'Qty (Previous)' in display_df.columns: format_dict['Qty (Previous)'] = '{:,.0f}'

            if 'Gross Profit (Current)' in display_df.columns: format_dict['Gross Profit (Current)'] = '{:,.0f}'
            if 'Gross Profit (%)' in display_df.columns: format_dict['Gross Profit (%)'] = '{:.2f}%'

            if 'Gross Margin (Current)' in display_df.columns: format_dict['Gross Margin (Current)'] = '{:,.0f}'
            if 'Gross Margin (%)' in display_df.columns: format_dict['Gross Margin (%)'] = '{:.2f}%'

            if 'Contribution Margin (Current)' in display_df.columns: format_dict[
                'Contribution Margin (Current)'] = '{:,.0f}'
            if 'Contribution Margin (%)' in display_df.columns: format_dict['Contribution Margin (%)'] = '{:.2f}%'
            if 'Qty Growth (%)' in display_df.columns: format_dict['Qty Growth (%)'] = '{:.2f}%'

            if 'Amount FG' in display_df.columns: format_dict['Amount FG'] = '{:,.0f}'
            if 'Amount Material' in display_df.columns: format_dict['Amount Material'] = '{:,.0f}'

            st.dataframe(
                display_df.style.format(format_dict, na_rep="0.00%"),
                use_container_width=True,
                hide_index=True
            )

            # ---------------------------------------------------------
            # PARETO ANALYSIS MODULE
            # ---------------------------------------------------------
            st.markdown("---")
            col_p_title, col_p_chart_metric, col_p_slider = st.columns([1.2, 1.4, 1.8])

            with col_p_title:
                st.subheader("📊 Pareto Analysis")

            with col_p_chart_metric:
                chart_margin_label = st.selectbox(
                    "Pilih Metric untuk Grafik:",
                    ["Gross Sales", "Gross Profit", "Gross Margin", "Contribution Margin"],
                    help="SKU diurutkan dari yang paling besar ke paling kecil berdasarkan metric ini."
                )

            with col_p_slider:
                pareto_threshold = st.slider(
                    "Pareto Threshold (% SKU)",
                    min_value=1.0,
                    max_value=100.0,
                    value=80.0,
                    step=1.0,
                    help="Berbasis JUMLAH SKU (bukan cumulative value). Contoh: 80% dari 100 SKU = 80 SKU teratas (diurutkan besar ke kecil) berdasarkan metric margin yang dipilih."
                )

            chart_margin_map = {
                "Gross Sales": "Gross Sales (Current)",
                "Gross Profit": "Gross Profit (Current)",
                "Gross Margin": "Gross Margin (Current)",
                "Contribution Margin": "Contribution Margin (Current)"
            }
            chart_margin_col = chart_margin_map[chart_margin_label]

            margin_cols_all = ['Gross Sales (Current)', 'Gross Profit (Current)', 'Gross Margin (Current)',
                               'Contribution Margin (Current)']
            agg_cols = [chart_margin_col]
            for c in ['Qty (Current)'] + margin_cols_all:
                if c in filtered_df.columns and c not in agg_cols:
                    agg_cols.append(c)

            if chart_margin_col in filtered_df.columns:
                df_pareto = filtered_df.groupby(name_col)[agg_cols].sum().reset_index()
                df_pareto = df_pareto.sort_values(by=chart_margin_col, ascending=False).reset_index(drop=True)

                if 'Dynamic 9-Box Category' in filtered_df.columns:
                    box_category_map = filtered_df.groupby(name_col)['Dynamic 9-Box Category'].agg(
                        lambda s: s.iloc[0] if s.nunique() == 1 else f"Mixed ({s.nunique()} box)"
                    )
                    df_pareto['9-Box Kwadran'] = df_pareto[name_col].map(box_category_map)

                if not df_pareto.empty:
                    total_val = df_pareto[chart_margin_col].sum()
                    if total_val != 0:
                        df_pareto['Cumulative %'] = (df_pareto[chart_margin_col].cumsum() / total_val) * 100
                    else:
                        df_pareto['Cumulative %'] = 0.0

                    total_all_skus = len(df_pareto)

                    if pareto_threshold >= 100:
                        cutoff_count = total_all_skus
                    else:
                        cutoff_count = max(1, round(total_all_skus * pareto_threshold / 100))
                    df_pareto_filtered = df_pareto.iloc[:cutoff_count].copy()

                    total_pareto_skus = len(df_pareto_filtered)

                    summary_sku = total_pareto_skus
                    summary_qty = df_pareto_filtered[
                        'Qty (Current)'].sum() if 'Qty (Current)' in df_pareto_filtered.columns else 0
                    summary_sales = df_pareto_filtered[
                        'Gross Sales (Current)'].sum() if 'Gross Sales (Current)' in df_pareto_filtered.columns else 0

                    pct_sku = (summary_sku / total_all_skus * 100) if total_all_skus > 0 else 0
                    pct_sales = (summary_sales / filtered_df[
                        'Gross Sales (Current)'].sum() * 100) if 'Gross Sales (Current)' in filtered_df.columns and \
                                                                 filtered_df['Gross Sales (Current)'].sum() > 0 else 0

                    m1, m2, m3 = st.columns(3)
                    m1.metric("Jumlah SKU", f"{summary_sku:,}", f"{pct_sku:.1f}% dari {total_all_skus:,} SKU")
                    m2.metric("Quantity", f"{summary_qty:,.0f}")
                    m3.metric("Gross Sales", f"Rp {summary_sales:,.0f}", f"{pct_sales:.1f}% dari total")

                    margin_label_map = {
                        'Gross Profit (Current)': 'Gross Profit',
                        'Gross Margin (Current)': 'Gross Margin',
                        'Contribution Margin (Current)': 'Contribution Margin'
                    }

                    mg1, mg2, mg3 = st.columns(3)

                    loop_margins = ['Gross Profit (Current)', 'Gross Margin (Current)', 'Contribution Margin (Current)']
                    for col_widget, margin_col in zip([mg1, mg2, mg3], loop_margins):
                        if margin_col not in df_pareto_filtered.columns:
                            continue
                        margin_sum = df_pareto_filtered[margin_col].sum()
                        grand_total = filtered_df[margin_col].sum() if margin_col in filtered_df.columns else 0
                        cumulative_pct = (margin_sum / grand_total * 100) if grand_total > 0 else 0
                        margin_ratio_pct = (margin_sum / summary_sales * 100) if summary_sales > 0 else 0
                        label = margin_label_map[margin_col]
                        if margin_col == chart_margin_col:
                            label = f"📊 {label} (di Grafik)"
                        margin_sum_bn = margin_sum / 1e9
                        col_widget.metric(
                            label,
                            f"Rp {margin_sum_bn:,.1f} M ({cumulative_pct:.1f}%)",
                            f"{margin_ratio_pct:.1f}% margin ratio"
                        )

                    fig_pareto = go.Figure()
                    
                    # Logika warna bar Pareto: Hitam untuk SKU Commonized, Merah untuk Minus, Biru untuk Normal
                    bar_colors = []
                    for i, row in df_pareto_filtered.iterrows():
                        npm = ""
                        if 'New Product Name' in filtered_df.columns:
                            matches = filtered_df[filtered_df[name_col] == row[name_col]]
                            if not matches.empty:
                                npm = str(matches.iloc[0]['New Product Name'])
                        
                        if npm.startswith('C-'):
                            bar_colors.append('#000000') # Warna Hitam Pekat untuk SKU Commonized
                        elif row[chart_margin_col] < 0:
                            bar_colors.append('#c03d32') # Merah
                        else:
                            bar_colors.append('#38bdf8') # Biru

                    fig_pareto.add_trace(go.Bar(
                        x=df_pareto_filtered[name_col],
                        y=df_pareto_filtered[chart_margin_col],
                        name=chart_margin_label,
                        marker_color=bar_colors,
                        yaxis='y1'
                    ))

                    fig_pareto.add_trace(go.Scatter(
                        x=df_pareto_filtered[name_col],
                        y=df_pareto_filtered['Cumulative %'],
                        name=f'Cumulative % ({chart_margin_label})',
                        marker_color='#f97316',
                        mode='lines+markers',
                        yaxis='y2'
                    ))

                    fig_pareto.update_layout(
                        title=f"Pareto Chart: {chart_margin_label} (Top {pareto_threshold}% SKU by count | {total_pareto_skus} dari {total_all_skus} SKUs, diurutkan besar ke kecil)",
                        hovermode="x unified",
                        height=550,
                        xaxis=dict(showticklabels=False, title=f"SKUs (Ranked by {chart_margin_label})"),
                        yaxis=dict(title=f"{chart_margin_label} (IDR)"),
                        yaxis2=dict(
                            title=f"Cumulative % {chart_margin_label}",
                            overlaying='y',
                            side='right',
                            showgrid=False
                        ),
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
                    )

                    st.plotly_chart(fig_pareto, use_container_width=True)

                    col_pt_dl, _ = st.columns([1, 3])
                    buffer_excel_pareto = io.BytesIO()
                    with pd.ExcelWriter(buffer_excel_pareto, engine='openpyxl') as writer:
                        df_pareto_filtered.to_excel(writer, index=False, sheet_name='Pareto_Data')

                    with col_pt_dl:
                        st.download_button(
                            label=f"📥 Download Top {pareto_threshold}% Data (Excel)",
                            data=buffer_excel_pareto.getvalue(),
                            file_name=f"Pareto_Analysis_{chart_margin_label.replace(' ', '_')}_{pareto_threshold}pctSKU.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            type="primary",
                            use_container_width=True
                        )

                    st.dataframe(
                        df_pareto_filtered.style.format({
                            'Gross Sales (Current)': '{:,.0f}',
                            'Qty (Current)': '{:,.0f}',
                            'Gross Profit (Current)': '{:,.0f}',
                            'Gross Margin (Current)': '{:,.0f}',
                            'Contribution Margin (Current)': '{:,.0f}',
                            'Cumulative %': '{:.2f}%'
                        }),
                        use_container_width=True,
                        hide_index=True
                    )

                    # =========================================================
                    # TAMBAHAN: 9-BOX DIAGRAM (MATRIX & BUBBLE) UNTUK HASIL PARETO
                    # =========================================================
                    st.markdown("---")
                    st.subheader(f"🧩 9-Box Summary untuk Top {pareto_threshold}% Pareto SKUs")
                    st.info(
                        "Visualisasi di bawah ini (Matrix Grid & Bubble Chart) menunjukkan sebaran 9-box khusus untuk SKU yang masuk dalam filter Pareto di atas.")

                    pareto_sku_list = df_pareto_filtered[name_col].tolist()
                    df_pareto_9box = filtered_df[filtered_df[name_col].isin(pareto_sku_list)].copy()

                    if not df_pareto_9box.empty:
                        tab_p_matrix, tab_p_bubble = st.tabs(["🧮 Matrix Grid", "🫧 Bubble Chart"])

                        with tab_p_matrix:
                            html_matrix_pareto = render_9box_summary_grid(
                                df_pareto_9box,
                                margin_selector,
                                margin_val_col,
                                y_low_thresh,
                                y_high_thresh,
                                x_axis_selector,
                                x_low_thresh,
                                x_high_thresh
                            )

                            st.download_button(
                                label="📥 Download Matrix Grid (.html)",
                                data=html_matrix_pareto,
                                file_name=f"Pareto_9Box_Matrix_{margin_selector.replace(' ', '_')}.html",
                                mime="text/html",
                                type="secondary"
                            )

                        with tab_p_bubble:
                            plot_df_p = df_pareto_9box.copy()

                            if bubble_size_selector in plot_df_p.columns:
                                plot_df_p['Bubble_Size'] = pd.to_numeric(plot_df_p[bubble_size_selector],
                                                                         errors='coerce').fillna(0).abs().replace(0, 1)
                            else:
                                plot_df_p['Bubble_Size'] = 10

                            box_counts_p = plot_df_p['Dynamic 9-Box Category'].value_counts().to_dict()
                            total_items_p = len(plot_df_p)

                            plot_df_p['X_Plot'] = plot_df_p[x_col].apply(apply_custom_x_scale)
                            plot_df_p['Y_Plot'] = plot_df_p[y_col]

                            # Override the color for Commonized SKUs to Black
                            plot_df_p['Plot_Color_Category'] = plot_df_p['Dynamic 9-Box Category']
                            if 'New Product Name' in plot_df_p.columns:
                                mask_commonized_p = plot_df_p['New Product Name'].fillna('').astype(str).str.startswith('C-')
                                plot_df_p.loc[mask_commonized_p, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

                            hover_data_dict_p = hover_data_dict.copy()

                            fig_p = px.scatter(
                                plot_df_p,
                                x='X_Plot',
                                y='Y_Plot',
                                size="Bubble_Size",
                                color="Plot_Color_Category",
                                color_discrete_map=color_discrete_map,
                                hover_name=name_col,
                                hover_data=hover_data_dict_p,
                                title=f"Pareto 9-Box Bubble Chart: {margin_selector} vs {x_axis_selector} (Total: {total_items_p} SKUs)",
                                size_max=60,
                                render_mode="svg"
                            )

                            fig_p.add_vline(x=apply_custom_x_scale(x_low_thresh), line_dash="dash", line_color="black",
                                            line_width=2, opacity=0.8,
                                            annotation_text=f" Low ({format_x_val_str(x_low_thresh)})",
                                            annotation_position="top left")
                            fig_p.add_vline(x=apply_custom_x_scale(x_high_thresh), line_dash="dash", line_color="black",
                                            line_width=2, opacity=0.8,
                                            annotation_text=f" High ({format_x_val_str(x_high_thresh)})",
                                            annotation_position="top right")
                            fig_p.add_hline(y=y_low_thresh, line_dash="dash", line_color="black", line_width=2,
                                            opacity=0.8, annotation_text=f"Low ({y_low_thresh:.1f}%)",
                                            annotation_position="bottom right")
                            fig_p.add_hline(y=y_high_thresh, line_dash="dash", line_color="black", line_width=2,
                                            opacity=0.8, annotation_text=f"High ({y_high_thresh:.1f}%)",
                                            annotation_position="top right")

                            fig_p.add_vline(x=apply_custom_x_scale(robust_avg_x), line_dash="dot", line_color="red",
                                            line_width=2, opacity=0.6,
                                            annotation_text=f" ← AVG X ({format_x_val_str(robust_avg_x)})",
                                            annotation_position="top right", annotation_font_color="red")
                            fig_p.add_hline(y=avg_margin_pct, line_dash="dot", line_color="red", line_width=2,
                                            opacity=0.6, annotation_text=f"AVG Y ({avg_margin_pct:.1f}%)",
                                            annotation_position="top right", annotation_font_color="red")

                            for box_name, coords in box_coords.items():
                                count = box_counts_p.get(box_name, 0)
                                if count > 0:
                                    is_b5 = (coords['id'] == 'B5')
                                    fig_p.add_annotation(
                                        x=coords['x'],
                                        y=coords['y'],
                                        ax=mapped_x_high + (mapped_plot_x_max - mapped_x_high) * 0.05 if is_b5 else None,
                                        ay=mid_y_low if is_b5 else None,
                                        axref="x" if is_b5 else None,
                                        ayref="y" if is_b5 else None,
                                        xref="x",
                                        yref="y",
                                        text=f"<span style='color:#1f2937;'><b>{coords['id']}</b></span><br><b>{count} SKUs</b>",
                                        showarrow=is_b5,
                                        arrowhead=2 if is_b5 else 0,
                                        arrowsize=1 if is_b5 else 1,
                                        arrowwidth=1 if is_b5 else 0.1,
                                        arrowcolor="#1f2937" if is_b5 else None,
                                        standoff=0,
                                        xanchor="left" if is_b5 else "center",
                                        font=dict(size=13, color="#374151"),
                                        bgcolor="rgba(255, 255, 255, 0.9)" if is_b5 else "rgba(255, 255, 255, 0.8)",
                                        bordercolor="rgba(15, 23, 42, 0.8)" if is_b5 else "rgba(15, 23, 42, 0.3)",
                                        borderwidth=2 if is_b5 else 1,
                                        borderpad=4
                                    )

                            fig_p.add_annotation(
                                x=0.98,
                                y=0.98,
                                xref="paper",
                                yref="paper",
                                text=f"<b>TOTAL:<br>{total_items_p} SKUs</b>",
                                showarrow=False,
                                font=dict(size=13, color="white"),
                                bgcolor="#374151",
                                bordercolor="black",
                                borderwidth=1,
                                borderpad=5
                            )

                            fig_p.update_xaxes(
                                range=[mapped_plot_x_min, mapped_plot_x_max],
                                tickvals=tickvals_x,
                                ticktext=ticktext_x,
                                title_text="LTM Quantity Growth Rate (%)" if x_axis_selector == "Qty Growth (%)" else "Gross Sales Absolute (IDR)"
                            )
                            fig_p.update_yaxes(range=[plot_y_min, plot_y_max], title_text=f"{margin_selector} (%)")

                            fig_p.update_layout(
                                height=550,
                                hovermode="closest",
                                showlegend=True,
                                legend=dict(
                                    title=None,
                                    orientation="h",
                                    yanchor="top",
                                    y=-0.15,
                                    xanchor="center",
                                    x=0.5
                                ),
                                margin=dict(t=50, b=50, l=50, r=50)
                            )

                            st.plotly_chart(fig_p, use_container_width=True, config={
                                'displayModeBar': True,
                                'modeBarButtonsToAdd': ['toImage'],
                                'displaylogo': False,
                                'toImageButtonOptions': {
                                    'format': 'png',
                                    'filename': f'Pareto_9Box_Bubble_{margin_selector.replace(" ", "_")}',
                                    'height': 800,
                                    'width': 1400,
                                    'scale': 2
                                }
                            })

                            buffer_html_p = io.StringIO()
                            fig_p.write_html(buffer_html_p, include_plotlyjs='cdn')
                            st.download_button(
                                label="📥 Download Bubble Chart (.html)",
                                data=buffer_html_p.getvalue(),
                                file_name=f"Pareto_9Box_Bubble_{margin_selector.replace(' ', '_')}.html",
                                mime="text/html",
                                type="secondary"
                            )
                    else:
                        st.warning("Tidak ada data yang valid untuk dirender ke dalam 9-Box Diagram.")

                else:
                    st.warning(f"No data found for {chart_margin_label} to generate a Pareto Chart.")
            else:
                st.error(f"Metric '{chart_margin_label}' not available in the dataset.")

    # ---------------------------------------------------------
    # TAB 2: BOX 3 INTERSECTION (DYNAMIC COMPARISON)
    # ---------------------------------------------------------
    with tab_b3_intersection:
        st.markdown("### 🎚️ Intersection Settings & Thresholds")
        st.info(
            "Dynamically compare any two margin metrics to identify SKUs that fall into 'Box 3' (Low Margin, Low Growth) for BOTH metrics simultaneously. Data is seamlessly derived from the Master P&L (Tab 1).")

        with st.form("intersection_settings_form"):
            col_x_int, col_y_int = st.columns(2)

            with col_x_int:
                metric_x = st.selectbox("Select Metric 1 (X-Axis):",
                                        ["Gross Profit", "Gross Margin", "Contribution Margin"], index=1)
                col_x_low, col_x_high = st.columns(2)
                with col_x_low:
                    x_low_med = st.number_input(f"{metric_x} - Low to Med Threshold (%)", value=25.0, step=1.0)
                with col_x_high:
                    x_med_high = st.number_input(f"{metric_x} - Med to High Threshold (%)", value=40.0, step=1.0)

            with col_y_int:
                metric_y = st.selectbox("Select Metric 2 (Y-Axis):",
                                        ["Gross Profit", "Gross Margin", "Contribution Margin"], index=2)
                col_y_low, col_y_high = st.columns(2)
                with col_y_low:
                    y_low_med = st.number_input(f"{metric_y} - Low to Med Threshold (%)", value=10.0, step=1.0)
                with col_y_high:
                    y_med_high = st.number_input(f"{metric_y} - Med to High Threshold (%)", value=25.0, step=1.0)

            run_intersect = st.form_submit_button("▶ RUN INTERSECTION", type="primary")

        col_x_pct = f"{metric_x} (%)"
        col_x_amt = f"{metric_x} (Current)"
        col_y_pct = f"{metric_y} (%)"
        col_y_amt = f"{metric_y} (Current)"

        box3_mask = (
                (filtered_df['Qty Growth (%)'] < 0) &
                (filtered_df[col_x_pct] < x_low_med) &
                (filtered_df[col_y_pct] < y_low_med)
        )

        df_intersect = filtered_df[box3_mask].copy()

        if 'Gross Sales (Current)' in df_intersect.columns:
            df_intersect['Bubble_Size'] = df_intersect['Gross Sales (Current)'].fillna(0).abs().replace(0, 1)
        else:
            df_intersect['Bubble_Size'] = 10

        df_local_intersect = df_intersect[df_intersect['Source_Sheet'] == 'LOCAL']
        df_export_intersect = df_intersect[df_intersect['Source_Sheet'] == 'EXPORT']

        sub_tab_local, sub_tab_export = st.tabs(["🏙️ LOCAL MARKET (Intersection)", "🌍 EXPORT MARKET (Intersection)"])

        base_display_cols = ['Product Name', 'Gross Sales (Current)', 'Qty (Current)', 'Qty Growth (%)',
                             col_x_amt, col_x_pct, col_y_amt, col_y_pct]
        display_cols = list(dict.fromkeys(base_display_cols))

        format_dict_intersect = {
            'Gross Sales (Current)': 'Rp {:,.0f}',
            'Qty (Current)': '{:,.0f}',
            'Qty Growth (%)': '{:.2f}%',
            col_x_amt: 'Rp {:,.0f}',
            col_x_pct: '{:.2f}%',
            col_y_amt: 'Rp {:,.0f}',
            col_y_pct: '{:.2f}%'
        }

        with sub_tab_local:
            if not df_local_intersect.empty:
                total_local_skus = len(df_local_intersect)
                st.info(
                    f"**Insight:** Found **{total_local_skus}** SKUs in the Local Market classified as 'Box 3' under both {metric_x} and {metric_y}.")

                # Terapkan warna hitam untuk Commonized SKU
                df_local_intersect['Plot_Color_Category'] = 'Normal SKU'
                if 'New Product Name' in df_local_intersect.columns:
                    mask_c = df_local_intersect['New Product Name'].fillna('').astype(str).str.startswith('C-')
                    df_local_intersect.loc[mask_c, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

                fig_local = px.scatter(
                    df_local_intersect, x=col_x_pct, y=col_y_pct,
                    size="Bubble_Size", color="Plot_Color_Category", hover_name="Product Name",
                    color_discrete_map={'Normal SKU': '#38bdf8', 'Commonized SKU (Black)': '#000000'},
                    title=f"Box 3 Intersection: {metric_x} vs {metric_y} (Local) | Total: {total_local_skus} SKUs",
                    size_max=50, render_mode="svg"
                )

                fig_local.add_vline(x=x_low_med, line_dash="dash", line_color="black", line_width=1, opacity=0.5,
                                    annotation_text=f"{metric_x} Low-Med ({x_low_med:g}%)")
                fig_local.add_vline(x=x_med_high, line_dash="dash", line_color="black", line_width=1, opacity=0.5)
                fig_local.add_hline(y=y_low_med, line_dash="dash", line_color="black", line_width=1, opacity=0.5,
                                    annotation_text=f"{metric_y} Low-Med ({y_low_med:g}%)")
                fig_local.add_hline(y=y_med_high, line_dash="dash", line_color="black", line_width=1, opacity=0.5)
                fig_local.add_vline(x=0, line_dash="dash", line_color="red", line_width=2, opacity=0.5)
                fig_local.add_hline(y=0, line_dash="dash", line_color="red", line_width=2, opacity=0.5)

                fig_local.update_layout(height=550, hovermode="closest", showlegend=True, margin=dict(t=50, b=50))
                st.plotly_chart(fig_local, use_container_width=True,
                                config={'displayModeBar': True, 'displaylogo': False})

                st.markdown("### 📋 Intersection Detail (Local)")

                valid_display_cols = [c for c in display_cols if c in df_local_intersect.columns]
                df_local_display = df_local_intersect[valid_display_cols].copy()

                buffer_local = io.BytesIO()
                with pd.ExcelWriter(buffer_local, engine='openpyxl') as writer:
                    df_local_display.to_excel(writer, index=False, sheet_name='Local_Intersection')

                col_dl_local, _ = st.columns([1, 3])
                with col_dl_local:
                    st.download_button("📥 Download Local Intersection Data (Excel)", data=buffer_local.getvalue(),
                                       file_name='Box3_Intersection_Local.xlsx',
                                       mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                                       type="primary", use_container_width=True)

                st.dataframe(df_local_display.style.format(format_dict_intersect), use_container_width=True,
                             hide_index=True)
            else:
                st.warning(
                    f"No SKUs found in the Local Market that fit the 'Box 3' criteria for both {metric_x} and {metric_y}.")

        with sub_tab_export:
            if not df_export_intersect.empty:
                total_export_skus = len(df_export_intersect)
                st.info(
                    f"**Insight:** Found **{total_export_skus}** SKUs in the Export Market classified as 'Box 3' under both {metric_x} and {metric_y}.")

                # Terapkan warna hitam untuk Commonized SKU
                df_export_intersect['Plot_Color_Category'] = 'Normal SKU'
                if 'New Product Name' in df_export_intersect.columns:
                    mask_c = df_export_intersect['New Product Name'].fillna('').astype(str).str.startswith('C-')
                    df_export_intersect.loc[mask_c, 'Plot_Color_Category'] = 'Commonized SKU (Black)'

                fig_export = px.scatter(
                    df_export_intersect, x=col_x_pct, y=col_y_pct,
                    size="Bubble_Size", color="Plot_Color_Category", hover_name="Product Name",
                    color_discrete_map={'Normal SKU': '#38bdf8', 'Commonized SKU (Black)': '#000000'},
                    title=f"Box 3 Intersection: {metric_x} vs {metric_y} (Export) | Total: {total_export_skus} SKUs",
                    size_max=50, render_mode="svg"
                )

                fig_export.add_vline(x=x_low_med, line_dash="dash", line_color="black", line_width=1, opacity=0.5,
                                     annotation_text=f"{metric_x} Low-Med ({x_low_med:g}%)")
                fig_export.add_vline(x=x_med_high, line_dash="dash", line_color="black", line_width=1, opacity=0.5)
                fig_export.add_hline(y=y_low_med, line_dash="dash", line_color="black", line_width=1, opacity=0.5,
                                     annotation_text=f"{metric_y} Low-Med ({y_low_med:g}%)")
                fig_export.add_hline(y=y_med_high, line_dash="dash", line_color="black", line_width=1, opacity=0.5)
                fig_export.add_vline(x=0, line_dash="dash", line_color="red", line_width=2, opacity=0.5)
                fig_export.add_hline(y=0, line_dash="dash", line_color="red", line_width=2, opacity=0.5)

                fig_export.update_layout(height=550, hovermode="closest", showlegend=True, margin=dict(t=50, b=50))
                st.plotly_chart(fig_export, use_container_width=True,
                                config={'displayModeBar': True, 'displaylogo': False})

                st.markdown("### 📋 Intersection Detail (Export)")

                export_cols = [
                                  'Country'] + base_display_cols if 'Country' in df_export_intersect.columns else base_display_cols
                export_cols = list(dict.fromkeys(export_cols))
                valid_export_cols = [c for c in export_cols if c in df_export_intersect.columns]

                df_export_display = df_export_intersect[valid_export_cols].copy()

                buffer_export = io.BytesIO()
                with pd.ExcelWriter(buffer_export, engine='openpyxl') as writer:
                    df_export_display.to_excel(writer, index=False, sheet_name='Export_Intersection')

                col_dl_export, _ = st.columns([1, 3])
                with col_dl_export:
                    st.download_button("📥 Download Export Intersection Data (Excel)", data=buffer_export.getvalue(),
                                       file_name='Box3_Intersection_Export.xlsx',
                                       mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                                       type="primary", use_container_width=True)

                st.dataframe(df_export_display.style.format(format_dict_intersect), use_container_width=True,
                             hide_index=True)
            else:
                st.warning(
                    f"No SKUs found in the Export Market that fit the 'Box 3' criteria for both {metric_x} and {metric_y}.")

    # ---------------------------------------------------------
    # TAB 3: RATIONALIZATION PROGRESS MONITOR
    # ---------------------------------------------------------
    with tab_progress:
        st.markdown("### 📉 SKU Rationalization Progress Monitor")

        # --- PERHITUNGAN DINAMIS DARI SUMBER DATA MENTAH (df_raw) ---
        # 1. Pastikan P1 dan P2 SELALU dari data UNCOMMON (Original Granular)
        df_uncommon = df_raw.copy()

        def get_sum(df, col):
            return df[col].sum() if col in df.columns else 0

        # Phase 1: Initial Base (All data without exception, Uncommon)
        p1_sku = len(df_uncommon)
        p1_sales = get_sum(df_uncommon, 'Gross Sales (Current)')
        p1_gm = get_sum(df_uncommon, 'Gross Margin (Current)')
        p1_cm = get_sum(df_uncommon, 'Contribution Margin (Current)')

        # Phase 2: First Optimization (Uncommon, Filter Exact Remark != Disc & Renewal)
        if 'Remark' in df_uncommon.columns:
            # Menggunakan exact match (isin) agar presisi membuang tepat 768 SKU tanpa menyentuh tipe Remark lainnya
            remark_clean = df_uncommon['Remark'].astype(str).str.strip().str.upper()
            mask_exc = remark_clean.isin(['DISC', 'RENEWAL'])
            df_p2 = df_uncommon[~mask_exc]
        else:
            df_p2 = df_uncommon.copy()

        p2_sku = len(df_p2)
        p2_sales = get_sum(df_p2, 'Gross Sales (Current)')
        p2_gm = get_sum(df_p2, 'Gross Margin (Current)')
        p2_cm = get_sum(df_p2, 'Contribution Margin (Current)')

        # Phase 3: Current Trim (Status == "Active" AND data harus COMMON/CONSOLIDATED)
        df_p3_raw = df_uncommon.copy()
        if 'Status' in df_p3_raw.columns:
            df_p3_raw = df_p3_raw[df_p3_raw['Status'].astype(str).str.strip().str.upper() == 'ACTIVE']

        # Lakukan commonization (Grouping Master SKU) KHUSUS untuk P3
        if 'New Code' in df_p3_raw.columns and 'New Product Name' in df_p3_raw.columns:
            df_p3_raw['New Code'] = df_p3_raw['New Code'].fillna('UNKNOWN').astype(str)
            df_p3_raw['New Product Name'] = df_p3_raw['New Product Name'].fillna('UNKNOWN').astype(str)

            group_cols_p3 = ['Source_Sheet', 'New Code', 'New Product Name']
            num_cols_p3 = df_p3_raw.select_dtypes(include=[np.number]).columns.tolist()
            sum_cols_p3 = [c for c in num_cols_p3 if not c.endswith('(%)')]

            agg_dict_p3 = {c: 'sum' for c in sum_cols_p3}
            df_p3_common = df_p3_raw.groupby(group_cols_p3, as_index=False).agg(agg_dict_p3)
        else:
            df_p3_common = df_p3_raw.copy()

        p3_sku = len(df_p3_common)
        p3_sales = get_sum(df_p3_common, 'Gross Sales (Current)')
        p3_gm = get_sum(df_p3_common, 'Gross Margin (Current)')
        p3_cm = get_sum(df_p3_common, 'Contribution Margin (Current)')

        # Membangun dataframe murni untuk rendering grafik
        progress_data_numeric = pd.DataFrame({
            'Phase': ['P1 (Initial Base)', 'P2 (First Optimization)', 'P3 (Current Trim)'],
            'Jumlah SKU': [p1_sku, p2_sku, p3_sku],
            'Gross Sales (IDR)': [p1_sales, p2_sales, p3_sales],
            'Gross Margin (IDR)': [p1_gm, p2_gm, p3_gm],
            'Contribution Margin (IDR)': [p1_cm, p2_cm, p3_cm]
        })

        # Helper untuk formatting tabel pakai koma sebagai ribuan
        def format_id_rupiah(val):
            return f"{int(val):,}"

        display_progress_df = progress_data_numeric.copy()
        display_progress_df['Jumlah SKU'] = display_progress_df['Jumlah SKU'].apply(format_id_rupiah)
        display_progress_df['Gross Sales (IDR)'] = display_progress_df['Gross Sales (IDR)'].apply(format_id_rupiah)
        display_progress_df['Gross Margin (IDR)'] = display_progress_df['Gross Margin (IDR)'].apply(format_id_rupiah)
        display_progress_df['Contribution Margin (IDR)'] = display_progress_df['Contribution Margin (IDR)'].apply(
            format_id_rupiah)

        st.dataframe(display_progress_df, use_container_width=True, hide_index=True)

        progress_metric = st.selectbox(
            "Select Tracking Metric for Y-Axis:",
            ["Jumlah SKU", "Gross Sales (IDR)", "Gross Margin (IDR)", "Contribution Margin (IDR)"]
        )

        # Distinct Corporate Palette (Corporate Navy, Executive Amber, Strategy Teal)
        executive_colors = ['#1e3a8a', '#d97706', '#0f766e']

        # Rendering the Chart using Numeric values (Plotly handles the formatting via separators)
        fig_prog = px.bar(
            progress_data_numeric,
            x='Phase',
            y=progress_metric,
            color='Phase',
            color_discrete_sequence=executive_colors,
            title=f"<b>RATIONALIZATION LIFECYCLE: {progress_metric.upper()}</b>"
        )

        if progress_metric == "Jumlah SKU":
            fig_prog.update_traces(
                texttemplate='<b>%{y:,.0f} SKUs</b>',
                textposition='outside',
                textfont=dict(size=16, color='#0f172a'),
                width=0.55, marker_line_width=0
            )
        else:
            fig_prog.update_traces(
                texttemplate='<b>Rp %{y:,.0f}</b>',
                textposition='outside',
                textfont=dict(size=16, color='#0f172a'),
                width=0.55, marker_line_width=0
            )

        fig_prog.update_layout(
            separators=",.",  # Force Plotly to use dot (.) as thousands separator and comma (,) as decimal
            plot_bgcolor='white',
            paper_bgcolor='white',
            showlegend=False,
            yaxis=dict(
                showgrid=True, gridcolor='#f1f5f9',
                showline=False, showticklabels=False, title=""
            ),
            xaxis=dict(
                showgrid=False, showline=True, linewidth=2, linecolor='#334155',
                title="", tickfont=dict(size=14, color='#334155', family="Arial")
            ),
            title_font=dict(size=22, color='#0f172a', family="Arial"),
            margin=dict(t=80, b=40, l=40, r=40),
            height=500
        )

        # Ensure Y-axis range is high enough so outside text doesn't get clipped
        max_y = progress_data_numeric[progress_metric].max()
        fig_prog.update_yaxes(range=[0, max_y * 1.2])

        st.plotly_chart(fig_prog, use_container_width=True)


if __name__ == "__main__":
    main()
